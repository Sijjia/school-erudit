# -*- coding: utf-8 -*-
# Сборка редактируемой презентации Bilim OS (12 слайдов) с ЧЁТКИМИ скринами.
# Стиль: светлый тёплый (оранж + индиго). Скрины кладутся отдельными картинками в полном разрешении.
# Запуск: python scripts/build_pptx.py
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.getcwd()
SHOTS = os.path.join(BASE, 'deck-shots')
OUT = os.path.join(BASE, 'school', 'deck', 'Bilim_OS_12_slides.pptx')

# палитра
INK   = RGBColor(0x1B,0x22,0x30)
MUT   = RGBColor(0x6C,0x76,0x86)
ACC   = RGBColor(0xF2,0x81,0x1D)
ACCD  = RGBColor(0xD9,0x6D,0x10)
IDG   = RGBColor(0x4F,0x46,0xE5)
BG    = RGBColor(0xFA,0xF6,0xF0)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LINE  = RGBColor(0xEA,0xDF,0xD0)
CARDBG= RGBColor(0xFF,0xFF,0xFF)
OKG   = RGBColor(0x1F,0xAA,0x6F)
FONT  = 'Segoe UI'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def bg_solid(s, rgb=BG):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb

def rect(s, l,t,w,h, fill=None, line=None, lw=1.0, rounded=False, radius=0.12, shadow=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(l),Inches(t),Inches(w),Inches(h))
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp

def grad_rect(s, l,t,w,h, c0=ACC, c1=IDG, angle=45):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l),Inches(t),Inches(w),Inches(h))
    shp.line.fill.background(); shp.shadow.inherit = False
    try:
        shp.fill.gradient()
        st = shp.fill.gradient_stops
        st[0].position = 0.0; st[0].color.rgb = c0
        st[1].position = 1.0; st[1].color.rgb = c1
        try: shp.fill.gradient_angle = angle
        except Exception: pass
    except Exception:
        shp.fill.solid(); shp.fill.fore_color.rgb = c0
    return shp

def text(s, l,t,w,h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6, line=1.12):
    """runs: list of paragraphs; each paragraph = list of (txt, size, bold, color, fontname)"""
    box = s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,para in enumerate(runs):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0)
        try: p.line_spacing = line
        except Exception: pass
        for (txt,size,bold,color,fn) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
            r.font.name = fn or FONT
    return box

def title(s, txt, num=None):
    rect(s, 0.7, 0.62, 0.10, 0.62, fill=ACC, rounded=True, radius=0.5)
    text(s, 0.95, 0.55, 11.0, 0.85, [[(txt, 32, True, INK, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    if num:
        text(s, 11.5, 0.55, 1.1, 0.6, [[(num, 13, True, MUT, FONT)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def footer(s, n):
    text(s, 0.7, 7.0, 4.0, 0.35, [[('BILIM OS', 10, True, ACC, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 9.0, 7.0, 3.6, 0.35, [[('Слайд %d из 12' % n, 10, False, MUT, FONT)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def feat_li(s, l, t, w, items):
    """буллеты-строки с оранжевой галочкой"""
    y = t
    for (mark, txt) in items:
        text(s, l, y, 0.4, 0.5, [[(mark, 15, True, ACC, FONT)]], anchor=MSO_ANCHOR.TOP)
        text(s, l+0.42, y, w-0.42, 0.8, [[(txt, 15, False, INK, FONT)]], anchor=MSO_ANCHOR.TOP, line=1.12)
        y += 0.62
    return y

def window(s, img, l, t, w, dots=True):
    """белая рамка-окно + картинка в полном разрешении"""
    im = Image.open(img); iw, ih = im.size; ratio = ih/iw
    bar = 0.32 if dots else 0.0
    pad = 0.11
    pic_w = w - pad*2
    pic_h = pic_w * ratio
    frame_h = bar + pad + pic_h + pad
    rect(s, l, t, w, frame_h, fill=WHITE, line=LINE, lw=1.0, rounded=True, radius=0.045)
    if dots:
        rect(s, l, t, w, bar+0.05, fill=RGBColor(0xF4,0xEE,0xE4), line=None, rounded=True, radius=0.05)
        for i,c in enumerate([RGBColor(0xFF,0x5F,0x57),RGBColor(0xFE,0xBC,0x2E),RGBColor(0x28,0xC8,0x40)]):
            d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l+0.22+i*0.22), Inches(t+0.11), Inches(0.13), Inches(0.13))
            d.fill.solid(); d.fill.fore_color.rgb=c; d.line.fill.background(); d.shadow.inherit=False
    s.shapes.add_picture(img, Inches(l+pad), Inches(t+bar+pad), width=Inches(pic_w))
    return frame_h

def phone(s, img, l, t, w, cap):
    im = Image.open(img); iw, ih = im.size; ratio = ih/iw
    pad = 0.07
    pic_w = w - pad*2; pic_h = pic_w*ratio
    fh = pic_h + pad*2
    rect(s, l, t, w, fh, fill=WHITE, line=LINE, lw=1.2, rounded=True, radius=0.10)
    s.shapes.add_picture(img, Inches(l+pad), Inches(t+pad), width=Inches(pic_w))
    text(s, l-0.3, t+fh+0.06, w+0.6, 0.3, [[(cap, 11, False, MUT, FONT)]], align=PP_ALIGN.CENTER)
    return fh

def card(s, l,t,w,h, icon, head, body, fill=None, idg=False):
    fillc = fill
    rect(s, l, t, w, h, fill=(fillc or WHITE), line=(None if fillc else LINE), lw=1.0, rounded=True, radius=0.07)
    icbg = RGBColor(0xEC,0xEB,0xFC) if idg else RGBColor(0xFD,0xF0,0xE2)
    iccol = IDG if idg else ACC
    if fillc: icbg = RGBColor(0xFF,0xFF,0xFF);
    ic = rect(s, l+0.32, t+0.30, 0.62, 0.62, fill=icbg, rounded=True, radius=0.28)
    tfi = ic.text_frame; tfi.word_wrap=False; tfi.margin_left=0; tfi.margin_right=0; tfi.margin_top=0; tfi.margin_bottom=0
    pi=tfi.paragraphs[0]; pi.alignment=PP_ALIGN.CENTER; ri=pi.add_run(); ri.text=icon; ri.font.size=Pt(20)
    hc = WHITE if fillc else INK
    bc = RGBColor(0xFF,0xFF,0xFF) if fillc else MUT
    text(s, l+0.32, t+1.10, w-0.6, 0.5, [[(head, 18, True, (WHITE if fillc else INK), FONT)]])
    text(s, l+0.32, t+1.65, w-0.6, h-1.8, [[(body, 13.5, False, bc, FONT)]], line=1.18)

# ============ СЛАЙД 1 — ТИТУЛ ============
s = slide(); grad_rect(s, 0,0,SW,SH, ACC, IDG, 45)
rect(s, 0.85, 1.7, 3.3, 0.5, fill=None, line=RGBColor(0xFF,0xFF,0xFF), lw=1.0, rounded=True, radius=0.5)
text(s, 0.85, 1.72, 3.3, 0.46, [[('ФОРУМ ДИРЕКТОРОВ ШКОЛ КР', 11, True, WHITE, FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.82, 2.4, 11.0, 1.7, [[('Bilim OS', 80, True, WHITE, FONT)]])
text(s, 0.85, 4.25, 10.5, 1.1,
     [[('Первая школьная ИИ-платформа в Кыргызстане. ', 21, False, WHITE, FONT)],
      [('От фиксации событий — к автоматическому действию.', 21, False, WHITE, FONT)]], line=1.2)
xx=0.85
for ch in ['ИИ-агенты 24/7','9 ролей','20+ модулей']:
    w = 0.5 + len(ch)*0.115
    rect(s, xx, 5.55, w, 0.5, fill=RGBColor(0xFF,0xFF,0xFF), line=None, rounded=True, radius=0.5)
    # полупрозрачность недоступна — делаем белый текст на лёгкой подложке
    cap = rect(s, xx, 5.55, w, 0.5, fill=None, line=RGBColor(0xFF,0xFF,0xFF), lw=1.0, rounded=True, radius=0.5)
    tf=cap.text_frame; tf.margin_top=0; tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text=ch; r.font.size=Pt(13); r.font.bold=True; r.font.color.rgb=ACCD
    xx += w + 0.25
text(s, 0.7, 7.0, 4.0, 0.35, [[('BILIM OS', 10, True, WHITE, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, 9.0, 7.0, 3.6, 0.35, [[('bilimos.kg', 10, False, WHITE, FONT)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# ============ СЛАЙД 2 — ПРОБЛЕМА ============
s = slide(); bg_solid(s)
title(s, 'Директор тонет в рутине')
probs = ['Бумажные журналы дублируют Excel-таблицы',
         'Родители узнают об оценках и пропусках с опозданием',
         'Отчёты собираются вручную неделями',
         'Управленческие решения принимаются вслепую, без аналитики']
y=2.0
for i,p in enumerate(probs):
    fillc = RGBColor(0xEC,0xEB,0xFC) if i%2 else RGBColor(0xFD,0xF0,0xE2)
    rect(s, 0.95, y, 11.4, 0.92, fill=WHITE, line=LINE, lw=1.0, rounded=True, radius=0.16)
    rect(s, 0.95, y, 0.14, 0.92, fill=(IDG if i%2 else ACC), rounded=True, radius=0.5)
    text(s, 1.35, y, 10.8, 0.92, [[(p, 17, False, INK, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.08
footer(s, 2)

# ============ СЛАЙД 3 — КОНЦЕПЦИЯ ============
s = slide(); bg_solid(s)
title(s, 'Событие → Действие. Без напоминаний')
# event card (indigo)
c1 = rect(s, 0.95, 2.0, 5.5, 4.2, fill=IDG, rounded=True, radius=0.05)
text(s, 1.35, 2.3, 4.8, 0.4, [[('1. СОБЫТИЕ — ВХОД', 12, True, WHITE, FONT)]])
text(s, 1.35, 2.75, 4.8, 0.6, [[('ИИ распознаёт инцидент', 20, True, WHITE, FONT)]])
ev=['Ученик получил низкую оценку','Накопились пропуски без причины','Успеваемость класса упала']
yy=3.6
for e in ev:
    text(s,1.35,yy,0.4,0.4,[[('●',13,True,WHITE,FONT)]]); text(s,1.75,yy,4.4,0.6,[[(e,15,False,WHITE,FONT)]]); yy+=0.62
# action card (orange)
rect(s, 6.85, 2.0, 5.5, 4.2, fill=ACC, rounded=True, radius=0.05)
text(s, 7.25, 2.3, 4.8, 0.4, [[('2. ДЕЙСТВИЕ — ВЫХОД', 12, True, WHITE, FONT)]])
text(s, 7.25, 2.75, 4.8, 0.6, [[('ИИ выполняет алгоритм', 20, True, WHITE, FONT)]])
ac=['Уведомление родителю в Telegram','Задача учителю — связаться с семьёй','Черновик отчёта директору']
yy=3.6
for a in ac:
    text(s,7.25,yy,0.4,0.4,[[('⚡',13,True,WHITE,FONT)]]); text(s,7.65,yy,4.4,0.6,[[(a,15,False,WHITE,FONT)]]); yy+=0.62
footer(s, 3)

# ============ СЛАЙД 4 — ПАНЕЛЬ ДИРЕКТОРА ============
s = slide(); bg_solid(s)
title(s, 'Панель управления директора')
text(s, 0.95, 1.95, 5.0, 0.5, [[('Вся школа на одном экране', 18, True, ACC, FONT)]])
feat_li(s, 0.95, 2.7, 5.0, [('✓','Живая сводка успеваемости и посещаемости'),
                            ('✓','Доступ ко всем модулям из сайдбара'),
                            ('✓','Мгновенные ИИ-оповещения об аномалиях')])
window(s, os.path.join(SHOTS,'02-dashboard-director.png'), 6.25, 2.05, 6.35)
footer(s, 4)

# ============ СЛАЙД 5 — ИИ-АГЕНТ ============
s = slide(); bg_solid(s)
title(s, 'Лента ИИ-агента')
# бейдж
b=rect(s, 5.0, 0.72, 1.95, 0.46, fill=RGBColor(0xDF,0xF6,0xEA), line=RGBColor(0xAE,0xE6,0xC9), lw=1.0, rounded=True, radius=0.5)
tf=b.text_frame; tf.margin_top=0; tf.margin_bottom=0; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
r=p.add_run(); r.text='УЖЕ РАБОТАЕТ'; r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=RGBColor(0x0E,0x7A,0x4F)
window(s, os.path.join(SHOTS,'07-agent-panel.png'), 0.7, 2.05, 6.55)
text(s, 7.55, 1.95, 5.0, 0.5, [[('Алерт · задача · совет · черновик', 17, True, ACC, FONT)]])
feat_li(s, 7.55, 2.7, 5.0, [('✓','Готовые задачи — не нужно держать в голове'),
                            ('✓','Черновики писем родителям пишутся сами'),
                            ('✓','Человек только жмёт «Согласовать и отправить»'),
                            ('✓','Свой инбокс агента у каждой роли')])
footer(s, 5)

# ============ СЛАЙД 6 — ЖУРНАЛ ============
s = slide(); bg_solid(s)
title(s, 'Электронный журнал')
b=rect(s, 4.5, 0.72, 1.95, 0.46, fill=RGBColor(0xDF,0xF6,0xEA), line=RGBColor(0xAE,0xE6,0xC9), lw=1.0, rounded=True, radius=0.5)
tf=b.text_frame; tf.margin_top=0; tf.margin_bottom=0; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
r=p.add_run(); r.text='УЖЕ РАБОТАЕТ'; r.font.size=Pt(11); r.font.bold=True; r.font.color.rgb=RGBColor(0x0E,0x7A,0x4F)
text(s, 0.95, 1.95, 5.0, 0.5, [[('Оценки в стиле EduPage', 18, True, ACC, FONT)]])
feat_li(s, 0.95, 2.7, 5.0, [('✓','Инлайн-выставление оценок в журнале класса'),
                            ('✓','Средний балл и дескрипторы — автоматически'),
                            ('✓','Низкая оценка сразу запускает ИИ-агента')])
window(s, os.path.join(SHOTS,'09-journal.png'), 6.25, 2.05, 6.35)
footer(s, 6)

# ============ СЛАЙД 7 — ИИ-ПРЕЗЕНТАЦИИ ============
s = slide(); bg_solid(s)
title(s, 'ИИ-презентации урока')
window(s, os.path.join(SHOTS,'10-presentations.png'), 0.7, 2.05, 6.55)
text(s, 7.55, 1.95, 5.0, 0.5, [[('Презентация за минуты', 18, True, ACC, FONT)]])
feat_li(s, 7.55, 2.7, 5.0, [('✓','Учитель задаёт тему — ИИ собирает презентацию'),
                            ('✓','Готово к выгрузке в PowerPoint'),
                            ('✓','Экономит часы подготовки к урокам')])
footer(s, 7)

# ============ СЛАЙД 8 — АНАЛИТИКА ============
s = slide(); bg_solid(s)
title(s, 'Интеллектуальный анализ школы')
text(s, 0.95, 1.95, 5.0, 0.5, [[('Прогнозирование успеваемости', 18, True, ACC, FONT)]])
feat_li(s, 0.95, 2.7, 5.0, [('⚠','Сигнал при спаде успеваемости'),
                            ('✎','Рекомендации по коррекции пробелов'),
                            ('✓','Готовые отчёты для директора')])
window(s, os.path.join(SHOTS,'03-analytics.png'), 6.25, 2.05, 6.35)
footer(s, 8)

# ============ СЛАЙД 9 — СМАРТФОН ============
s = slide(); bg_solid(s)
title(s, 'Школа в смартфоне')
text(s, 0.95, 1.85, 11.4, 0.5, [[('Родитель и ученик видят дневник, ДЗ и расписание прямо в телефоне — без звонков в школу.', 15, False, MUT, FONT)]])
phs=[('16-diary-parent-mobile.png','Дневник · родитель'),
     ('17-homework-student-mobile.png','Домашние задания · ученик'),
     ('18-schedule-student-mobile.png','Расписание · ученик')]
px=2.96
for img,cap in phs:
    phone(s, os.path.join(SHOTS,img), px, 2.45, 2.0, cap)
    px += 2.7
footer(s, 9)

# ============ СЛАЙД 10 — БЕЗОПАСНОСТЬ ============
s = slide(); bg_solid(s)
title(s, 'Данные под защитой')
card(s, 0.95, 2.2, 3.65, 3.4, '🔐', '9 ролей доступа', 'Жёсткое разграничение: каждый видит только своё. Ученик и родитель — только свои оценки.')
card(s, 4.85, 2.2, 3.65, 3.4, '📜', 'Журнал аудита', 'Кто, когда и что изменил — фиксируется. Корректность проверена автотестами.', idg=True)
card(s, 8.75, 2.2, 3.65, 3.4, '🛡️', 'Приватная обработка ИИ', 'Данные детей — по принципам COPPA/GDPR. Облако и регулярные бэкапы.')
footer(s, 10)

# ============ СЛАЙД 11 — СРАВНЕНИЕ ============
s = slide(); bg_solid(s)
title(s, 'Сравнение систем управления')
rows = [('Функция','Обычные системы','Bilim OS (ИИ)'),
        ('Контроль посещаемости','Заполнение журнала вручную','Автоматическая фиксация'),
        ('Оповещение родителей','Классный руководитель звонит','Мгновенно в Telegram'),
        ('Реакция на двойку','Замечают через недели','ИИ-агент в ту же секунду'),
        ('Подготовка урока','Часы ручной работы','Презентация за минуты')]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.95), Inches(2.05), Inches(11.4), Inches(4.2)).table
tbl.columns[0].width=Inches(3.6); tbl.columns[1].width=Inches(4.2); tbl.columns[2].width=Inches(3.6)
for ri,row in enumerate(rows):
    for ci,val in enumerate(row):
        cell = tbl.cell(ri,ci); cell.text=''
        p=cell.text_frame.paragraphs[0]; r=p.add_run(); r.text=val
        r.font.name=FONT; r.font.size=Pt(14.5)
        cell.vertical_anchor=MSO_ANCHOR.MIDDLE
        cell.margin_left=Inches(0.18); cell.margin_top=Inches(0.06); cell.margin_bottom=Inches(0.06)
        if ri==0:
            cell.fill.solid(); cell.fill.fore_color.rgb=RGBColor(0xF6,0xF0,0xE6)
            r.font.bold=True; r.font.color.rgb=(ACC if ci==2 else MUT); r.font.size=Pt(12.5)
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb=WHITE
            if ci==0: r.font.bold=True; r.font.color.rgb=INK
            elif ci==1: r.font.color.rgb=MUT
            else:
                r.font.bold=True; r.font.color.rgb=OKG
                r.text='✓ '+val
footer(s, 11)

# ============ СЛАЙД 12 — CTA ============
s = slide(); grad_rect(s, 0,0,SW,SH, ACC, IDG, 45)
rect(s, 0.85, 1.55, 2.7, 0.5, fill=None, line=RGBColor(0xFF,0xFF,0xFF), lw=1.0, rounded=True, radius=0.5)
text(s, 0.85, 1.57, 2.7, 0.46, [[('СОБЫТИЕ → ДЕЙСТВИЕ', 11, True, WHITE, FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.82, 2.3, 11.5, 2.0, [[('Станьте пионером ИИ', 52, True, WHITE, FONT)],
                               [('в школах Кыргызстана', 52, True, WHITE, FONT)]], line=1.02)
text(s, 0.85, 4.75, 10.8, 0.9, [[('Приглашаем прогрессивные школы в пилот Bilim OS — и зафиксируйте условия раннего партнёра.', 18, False, WHITE, FONT)]], line=1.25)
text(s, 0.85, 5.85, 11.8, 0.6, [[('bilimos.kg     ·     Askarova.datkayim@gmail.com     ·     +996 700 144 03', 17, True, WHITE, FONT)]])
text(s, 0.7, 7.0, 4.0, 0.35, [[('BILIM OS', 10, True, WHITE, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, 9.0, 7.0, 3.6, 0.35, [[('Слайд 12 из 12', 10, False, WHITE, FONT)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print('OK ->', OUT)
