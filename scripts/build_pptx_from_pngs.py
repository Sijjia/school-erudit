# -*- coding: utf-8 -*-
# PPTX из отрендеренных слайдов (school/deck/slides-png, 2x ретина).
# Каждый слайд — полнокадровая картинка 16:9. Для показа/печати; правки — в deck.html.
# Запуск: python scripts/build_pptx_from_pngs.py
import os, glob, sys
from pptx import Presentation
from pptx.util import Inches

BASE = os.getcwd()
NAME = sys.argv[1] if len(sys.argv) > 1 else ''  # '' = полная дека, '12' = короткая
PNG_DIR = 'slides-png' + (f'-{NAME}' if NAME else '')
PNGS = sorted(glob.glob(os.path.join(BASE, 'school', 'deck', PNG_DIR, 'slide-*.png')))
OUT = os.path.join(BASE, 'school', 'deck', f'Bilim_OS_deck{"_" + NAME if NAME else ""}.pptx')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

for p in PNGS:
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(OUT)
print(f'OK {len(PNGS)} slides -> {OUT}')
